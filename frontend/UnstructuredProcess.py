import os
import docx
import pptx
import PyPDF2
import pandas as pd
import streamlit as st
import re
import requests
import json


def add_title_step_1():
    """
    Adding title for step 1 - specifying file type
    """
    st.write("------------------------------------\n------------------------------------")
    st.markdown(f"### Step 1: Specifying the root directory")


def add_title_step_2(state):
    """
    Adding title for step 2 - begin chatting
    """

    if state.method_two_df is not None:
        st.write("------------------------------------\n------------------------------------")
        st.markdown(f"### Step 2: Information extraction via user query")


def get_root_dir(state, status):

    status.info("Please provide the root directory")

    state.root_dir = st.text_input("Please provide the root directory to be searched through:")

    return state, status


def get_filetypes(state, status):

    if state.root_dir is not None or state.root_dir == "":

        status.info("Please provide the file types to be searched for")

        state.file_types_options = st.multiselect(
            'Please select the file type(s) to be searched for:',
            ['.pptx', 
            '.pdf', 
            '.docx']
            )

    return state, status


def begin_file_dataframe_gen_process(state, status):


    if (state.file_types_options is not None) or \
        (len(state.file_types_options) != 0):

        status.info("Waiting to begin data collation into dataframe process...")
        st.write("") #for spacing only
        if st.button("Begin file data collation"):

            state.begin_file_data_collation = True

    return state, status


def createlist(root_dir_path, file_extensions):
    """
    Returns a list of all files with specified extensions
    found within the specified root directory path
    """
    path = root_dir_path
    list_files = []
    count = 0

    for root, dirs, files in os.walk(path, topdown=False):
        for file in files:
            if file.endswith(tuple(file_extensions)):
                count += 1
                file_list = os.path.join(root, file)
                file_list = os.path.relpath(file_list, root_dir_path)
                list_files.append(file_list)

    return list_files


def process_docx(file_path):
    """
    Process .docx files and combine the text content
    """
    doc = docx.Document(file_path)
    combined_text = ""

    for paragraph in doc.paragraphs:
        combined_text += paragraph.text + " "

    return combined_text


def process_pptx(file_path):
    """
    Process .pptx files and combine the text content
    """
    prs = pptx.Presentation(file_path)
    combined_text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    combined_text += run.text + " "

    return combined_text


def process_pdf(file_path):
    """
    Process .pdf files and combine the text content
    """
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        combined_text = ""

        for page in reader.pages:
            combined_text += page.extract_text()

    return combined_text


def process_files_into_df(root_dir_path, list_files):
    """
    Processes the identified files and saves content
    to a dictionary for later search entity matching
    """
    dict_data = {}

    for file in list_files:
        file_path = os.path.join(root_dir_path, file)
        file_extension = os.path.splitext(file_path)[1].lower()

        if file_extension == '.docx':
            combined_text = process_docx(file_path)
        elif file_extension == '.pptx':
            combined_text = process_pptx(file_path)
        elif file_extension == '.pdf':
            combined_text = process_pdf(file_path)
        else:
            continue

        # Add a space after each sentence
        combined_text = re.sub(r'(?<=[.!?])', ' ', combined_text)
        dict_data[file] = combined_text.rstrip()

    # Convert dictionary to DataFrame
    df = pd.DataFrame(dict_data.items(), columns=['Relative File Path', 'Information'])
    
    return df


def create_df_from_files_rootdir(state, status):

    if state.begin_file_data_collation is True and state.method_two_df is None:

        list_files = createlist(state.root_dir, state.file_types_options)
        state.method_two_df = process_files_into_df(state.root_dir, list_files)

        status.info("Dataframe generated!")

    return state, status


def display_df_on_app(state, status):

    if state.method_two_df is not None:

        status.info("Displaying generated dataframe...")
        st.write("")
        st.write(state.method_two_df)

    return status


def save_df_locally(state):

    if state.method_two_df is not None:

        state.method_two_df.to_csv("data_method_2/data_method_two_saved.csv", index=False)


def input_query_text_area(state, status):

    if state.method_two_df is not None:

        status.info("Please enter your query...")
        col1, _ = st.columns([0.7, 0.3])

        with col1:

            state.query_method_two = st.text_area("Please enter your query regarding the dataframe:")
        
        if st.button("Save query"):
            state.save_query_two = True

    return state, status


def generate_proper_prompt(state):

    if state.save_query_two:
        state.query_boundary_condition = "Do not use your foundation knowledge. Only answer questions based on information provided in the context."
        state.prompt_method_two = state.query_boundary_condition + " " + state.query_method_two
        state.que_response_gen_two = True
        state.response_method_two = None
    return state


def call_api(state, path_to_df: str, prompt: str):

    url = "http://localhost:8001/process_text_v2/"
    parameters = {"path_to_df": path_to_df, "prompt": prompt}
    
    response_two = requests.get(url, params=parameters)
    result_two = response_two.json()["processed_text"]

    state.response_method_two = result_two[len("Output: "):]

    return state


def generate_response(state, status):

    path_to_df2 = "data_method_2/data_method_two_saved.csv"

    if state.que_response_gen_two is True:
        status.info("Ready for response generation!")

        st.write("----------------------------------------------------")
        if st.button("Ask query"):
            # on_click=call_api, kwargs={"state": state, "path_to_df": path_to_df2, "prompt": state.prompt_method_two}
            with st.spinner("Generating response..."):
                state = call_api(state, path_to_df=path_to_df2, prompt=state.prompt_method_two)
    
    return state, status


def display_response(state, status):
    if state.response_method_two is not None:
        status.info("Dislaying response")
        _,  col2 = st.columns([0.3, 0.7])
        with col2:
            st.caption("Response:")
            st.success(state.response_method_two, icon="🤖")


def unstructured_pipeline(state, status):

    add_title_step_1()
    state, status = get_root_dir(state, status)
    state, status = get_filetypes(state, status)
    state, status = begin_file_dataframe_gen_process(state, status)
    state, status = create_df_from_files_rootdir(state, status)
    status = display_df_on_app(state, status)
    save_df_locally(state)
    add_title_step_2(state)
    state, status = input_query_text_area(state, status)
    state = generate_proper_prompt(state)
    state, status = generate_response(state, status)
    display_response(state, status)

    return state, status
