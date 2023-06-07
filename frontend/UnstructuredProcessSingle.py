import os
import docx
import pptx
import PyPDF2
import pandas as pd
import streamlit as st
import re
import requests
import json
import nltk
nltk.download('punkt')  # Download the required NLTK data
from nltk.tokenize import sent_tokenize


from frontend.streamlit_state_handling import *


def add_title_step_1():
    """
    Adding title for step 1 - specifying file type
    """
    st.write("------------------------------------\n------------------------------------")
    st.markdown(f"### Step 1: Specifying the filetype")
    st.markdown("##### Please upload file to be processed.")


def add_title_step_2(state):
    """
    Adding title for step 2 - begin chatting
    """

    if state.file_3 is not None:
        st.write("------------------------------------\n------------------------------------")
        st.markdown(f"### Step 2: Information extraction via user query")


def file_loading(state):
    """
    Uploading and processing jpg/png/tif file
    """

    # file is uploaded
    uploaded_file = st.file_uploader(label='Upload file of type:', type=['.docx', '.pdf', '.pptx'])
    
    if uploaded_file is None:

        state.file_3 = None
        state = state_upload_file_3_none(state)

    else:

        # file converted to numpy array for later processing

        state.file_3 = uploaded_file 
        state.file_path_3 = state.file_3.name

    return state


def process_docx(file_path):
    """
    Process .docx files and combine the text content
    """
    doc = docx.Document(file_path)
    combined_text = ""

    for paragraph in doc.paragraphs:
        combined_text += paragraph.text + " "

    return combined_text


def process_pptx(file_path):
    """
    Process .pptx files and combine the text content
    """
    prs = pptx.Presentation(file_path)
    combined_text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    combined_text += run.text + " "

    return combined_text


def process_pdf(file_path):
    """
    Process .pdf files and combine the text content
    """
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        combined_text = ""

        for page in reader.pages:
            combined_text += page.extract_text()

    return combined_text


def process_file_into_df_old(file_path):
    """
    Processes the identified file and saves content
    to a str object for later search
    """

    if file_path.endswith('.docx'):
        combined_text = process_docx(file_path)
    elif file_path.endswith('.pptx'):
        combined_text = process_pptx(file_path)
    elif file_path.endswith('.pdf'):
        combined_text = process_pdf(file_path)

    # Add a space after each sentence
    combined_text = re.sub(r'(?<=[.!?])', ' ', combined_text)
    
    return combined_text


def process_file_into_df_old_2(file_path):
    """
    Processes the identified file and saves content
    to a str object for later search
    """

    if file_path.endswith('.docx'):
        combined_text = process_docx(file_path)
    elif file_path.endswith('.pptx'):
        combined_text = process_pptx(file_path)
    elif file_path.endswith('.pdf'):
        combined_text = process_pdf(file_path)

    # Add a space after each sentence
    combined_text = re.sub(r'(?<=[.!?])', ' ', combined_text)
    
    sentences = re.split(r'(?<=[.!?])\s', combined_text)
    sentence_chunks = [sentences[i:i+5] for i in range(0, len(sentences), 5)]

    flat_sentences = [sentence for chunk in sentence_chunks for sentence in chunk]

    dict_data = {'Content': flat_sentences}
    df = pd.DataFrame(dict_data)
    
    return df


def process_file_into_df_old_3(file_path):
    """
    Processes the identified file and saves content to a str object for later search
    """

    if file_path.endswith('.docx'):
        combined_text = process_docx(file_path)
    elif file_path.endswith('.pptx'):
        combined_text = process_pptx(file_path)
    elif file_path.endswith('.pdf'):
        combined_text = process_pdf(file_path)

    # Add a space after each sentence
    combined_text = re.sub(r'(?<=[.!?])', ' ', combined_text)

    sentences = re.split(r'(?<=[.!?])\s', combined_text)

    # Remove any empty sentences
    sentences = [sentence for sentence in sentences if sentence.strip()]

    sentence_chunks = [sentences[i:i+5] for i in range(0, len(sentences), 5)]

    # Pad the last chunk with empty strings if necessary
    last_chunk = sentence_chunks[-1]
    if len(last_chunk) < 5:
        last_chunk += [''] * (5 - len(last_chunk))
        sentence_chunks[-1] = last_chunk

    flat_sentences = [sentence for chunk in sentence_chunks for sentence in chunk]

    dict_data = {'Content': flat_sentences}
    df = pd.DataFrame(dict_data)

    return df, combined_text



def process_file_into_df(file_path):
    """
    Processes the identified file and saves content to a str object for later search
    """

    if file_path.endswith('.docx'):
        combined_text = process_docx(file_path)
    elif file_path.endswith('.pptx'):
        combined_text = process_pptx(file_path)
    elif file_path.endswith('.pdf'):
        combined_text = process_pdf(file_path)

    # Add a space after each sentence
    combined_text = re.sub(r'(?<=[.!?])', ' ', combined_text)

    combined_text = combined_text.replace('\n', '')  # Remove newline symbols

    sentences = sent_tokenize(combined_text)

    # Group sentences into chunks of 5
    sentence_chunks = [sentences[i:i+5] for i in range(0, len(sentences), 5)]

    # Create a list of rows with 5 sentences each
    rows = []
    for chunk in sentence_chunks:
        if len(chunk) < 5:
            # Pad with empty sentences if the chunk has fewer than 5 sentences
            chunk += [''] * (5 - len(chunk))
        rows.append(' '.join(chunk))

    # Create DataFrame with 'sentences' as the column name
    df = pd.DataFrame(rows, columns=['Content'])

    return df, combined_text


def display_text_snippet_on_app(state):

    if state.file_path_3 is not None:
        _, combined_text = process_file_into_df(state.file_path_3)
        first_150_char = "**File Content Snippet**: \n\n"+combined_text[:700]+" (continued)..."
        st.info(first_150_char)


def input_query_text_area(state, status):

    if state.file_3 is not None:

        status.info("Please enter your query...")
        col1, _ = st.columns([0.7, 0.3])

        with col1:

            state.query_method_three = st.text_area("Please enter your query regarding uploaded file:")
        
        if st.button("Save query"):
            state.save_query_three = True

    return state, status


def generate_proper_prompt(state):

    if state.save_query_three:
        state.query_boundary_condition = "Do not use your foundation knowledge. Only answer questions based on information provided in the context."
        state.prompt_method_three = state.query_boundary_condition + " " + state.query_method_three
        state.que_response_gen_three = True
        state.response_method_three = None
    return state


def call_api(state, path_to_df: str, prompt: str):

    url = "http://localhost:8002/process_text_v3/"
    parameters = {"path_to_df": path_to_df, "prompt": prompt}
    
    response_three = requests.get(url, params=parameters)
    result_three = response_three.json()["processed_text"]

    state.response_method_three = result_three[len("Output: "):]

    return state


def generate_response(state, status):

    if state.que_response_gen_three is True:
        status.info("Ready for response generation!")

        st.write("----------------------------------------------------")
        if st.button("Ask query"):
            # on_click=call_api, kwargs={"state": state, "path_to_df": path_to_df2, "prompt": state.prompt_method_three}
            with st.spinner("Generating response..."):
                state = call_api(state, path_to_df=state.file_path_3, prompt=state.prompt_method_three)
    
    return state, status


def display_response(state, status):
    if state.response_method_three is not None:
        status.info("Dislaying response")
        _,  col2 = st.columns([0.3, 0.7])
        with col2:
            st.caption("Response:")
            st.success(state.response_method_three, icon="🤖")


def unstructured_pipeline_single(state, status):

    add_title_step_1()
    state = file_loading(state)
    display_text_snippet_on_app(state)
    add_title_step_2(state)
    state, status = input_query_text_area(state, status)
    state = generate_proper_prompt(state)
    state, status = generate_response(state, status)
    display_response(state, status)

    return state, status
